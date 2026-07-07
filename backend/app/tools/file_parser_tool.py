# file_parser_tool.py
# Purpose: Extract clean text payloads from uploaded document files.
# Responsibilities:
#   - Parse PDF, DOCX, PPTX, and TXT extensions using proper format decoders
#   - Handle exceptions related to corrupt files, password protection, and encoding errors
# DO NOT: Store parsed files on disk (process them directly from memory/file-streams).
# DO NOT: Run embedding operations directly in this file.

import io
import logging
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

from app.core.exceptions import FileParsingException

logger = logging.getLogger(__name__)


def parse_file(file_content: bytes, filename: str) -> str:
    """
    Parse uploaded files and return raw text content.

    Args:
        file_content: The file bytes.
        filename: Name of the file containing the extension.

    Returns:
        str: Raw text extracted.
    """
    # 1. Sanitize filename to prevent traversal or shell escapes
    import os
    basename = os.path.basename(filename)
    sanitized_name = "".join(c for c in basename if c.isalnum() or c in (".", "-", "_"))
    if not sanitized_name:
        sanitized_name = "document.txt"

    # 2. Validate file size matches configured limits
    from app.core.config import settings
    max_size_bytes = settings.MAX_FILE_SIZE_MB * 1024 * 1024
    if len(file_content) > max_size_bytes:
        raise FileParsingException(
            f"File size exceeds configured limit of {settings.MAX_FILE_SIZE_MB}MB."
        )

    ext = sanitized_name.lower()
    try:
        if ext.endswith(".pdf"):
            return _parse_pdf(file_content)
        elif ext.endswith(".docx"):
            return _parse_docx(file_content)
        elif ext.endswith(".pptx"):
            return _parse_pptx(file_content)
        elif ext.endswith(".txt"):
            return file_content.decode("utf-8", errors="ignore")
        else:
            raise FileParsingException(f"Unsupported file format for: {filename}")
    except Exception as e:
        logger.error(f"Error parsing file {filename}: {e}")
        if isinstance(e, FileParsingException):
            raise
        raise FileParsingException(f"Failed to parse {filename}: {str(e)}")


def _parse_pdf(content: bytes) -> str:
    """Extract text from PDF file bytes."""
    pdf_file = io.BytesIO(content)
    reader = PdfReader(pdf_file)
    text = []
    for page in reader.pages:
        extracted = page.extract_text()
        if extracted:
            text.append(extracted)
    return "\n".join(text)


def _parse_docx(content: bytes) -> str:
    """Extract text from DOCX paragraphs."""
    docx_file = io.BytesIO(content)
    doc = Document(docx_file)
    return "\n".join([p.text for p in doc.paragraphs])


def _parse_pptx(content: bytes) -> str:
    """Extract text from PPTX slide shapes."""
    pptx_file = io.BytesIO(content)
    prs = Presentation(pptx_file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text.append(shape.text)
    return "\n".join(text)
